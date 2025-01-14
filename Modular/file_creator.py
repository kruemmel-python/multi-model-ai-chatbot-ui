from openpyxl import Workbook
from docx import Document
from fpdf import FPDF
from pptx import Presentation
import csv
from api_client import api_client
from config import MISTRAL_CHAT_MODEL, DEFAULT_OLLAMA_MODEL
import subprocess
import re
import os

class FileCreator:
    """
    Eine Klasse zur Erstellung verschiedener Dateitypen (Excel, Word, PDF, PowerPoint, CSV) 
    mit Inhalten, die von KI-Modellen generiert werden.
    """

    def __init__(self):
        """Initialisiert eine Instanz der FileCreator-Klasse."""
        pass

    def generate_content_with_model(self, model_name: str, user_prompt: str) -> str:
        """
        Generiert den Inhalt basierend auf einem ausgewählten KI-Modell.

        Args:
            model_name (str): Der Name des zu verwendenden Modells (z. B. "Mistral", "Gemini", "Ollama").
            user_prompt (str): Der Benutzer-Input, der an das Modell übergeben wird.

        Returns:
            str: Der generierte Inhalt als Zeichenkette.
        """
        if model_name == "Mistral":
            # Anfrage an das Mistral-Modell
            response = api_client.mistral_client.chat.complete(
                model=MISTRAL_CHAT_MODEL,
                messages=[{"role": "user", "content": user_prompt}]
            )
            return response.choices[0].message.content.strip()
        elif model_name == "Gemini":
            # Anfrage an das Gemini-Modell
            response = api_client.gemini_model.generate_content([user_prompt])
            return response.text.strip()
        elif model_name == "Ollama":
            # Anfrage an das Ollama-Modell über einen Subprozess
            process = subprocess.Popen(
                ["ollama", "run", DEFAULT_OLLAMA_MODEL],
                stdin=subprocess.PIPE,
                stdout=subprocess.PIPE,
                text=True
            )
            process.stdin.write(user_prompt + "\n")
            process.stdin.close()
            output = process.stdout.read()
            process.stdout.close()
            return self.clean_output(output).strip()
        else:
            return "Modell nicht verfügbar oder unbekannt."

    def clean_output(self, output: str) -> str:
        """
        Entfernt Steuerzeichen aus der Ausgabe des Ollama-Modells.

        Args:
            output (str): Die ursprüngliche Ausgabe.

        Returns:
            str: Die bereinigte Ausgabe.
        """
        # Entfernen von ANSI-Steuerzeichen und anderen unerwünschten Mustern.
        cleaned_output = re.sub(r'(?:\x1B[@-_]|[\x1B\x9B][0-?]*[ -/]*[@-~])', '', output)
        cleaned_output = re.sub(r'\?\d+[lh]', '', cleaned_output)
        cleaned_output = re.sub(r'[\u2800-\u28FF]', '', cleaned_output)
        cleaned_output = re.sub(r'\r', '', cleaned_output)
        cleaned_output = re.sub(r'2K1G ?(?:2K1G)*!?', '', cleaned_output)
        return cleaned_output

    def create_excel_with_ai(self, user_prompt: str, sheets: int = 1) -> str:
        """
        Erstellt eine Excel-Datei mit von KI generiertem Inhalt.

        Args:
            user_prompt (str): Der Benutzer-Input, der für die Generierung verwendet wird.
            sheets (int): Die Anzahl der zu erstellenden Tabellenblätter (Standard: 1).

        Returns:
            str: Der Dateipfad der erstellten Excel-Datei.
        """
        content = self.generate_content_with_model("Gemini", user_prompt)
        workbook = Workbook()
        for i in range(sheets):
            # Erstellung oder Auswahl eines Tabellenblatts.
            sheet = workbook.create_sheet(title=f"Tabelle{i+1}") if i > 0 else workbook.active
            rows = content.split("\n")
            for row_index, row in enumerate(rows, start=1):
                columns = row.split(",")  # Annahme: Inhalte sind durch Kommas getrennt.
                for col_index, value in enumerate(columns, start=1):
                    sheet.cell(row=row_index, column=col_index, value=value.strip())
        file_path = "erstellte_tabelle.xlsx"
        workbook.save(file_path)
        os.startfile(file_path)  # Datei automatisch öffnen.
        return file_path

    def create_word_with_ai(self, user_prompt: str) -> str:
        """
        Erstellt eine Word-Datei mit von KI generiertem Inhalt.

        Args:
            user_prompt (str): Der Benutzer-Input für die Generierung.

        Returns:
            str: Der Dateipfad der erstellten Word-Datei.
        """
        content = self.generate_content_with_model("Gemini", user_prompt)
        doc = Document()
        doc.add_paragraph(content)
        file_path = "erstelltes_dokument.docx"
        doc.save(file_path)
        os.startfile(file_path)  # Datei automatisch öffnen.
        return file_path

    def create_pdf_with_ai(self, user_prompt: str) -> str:
        """
        Erstellt eine PDF-Datei mit von KI generiertem Inhalt.

        Args:
            user_prompt (str): Der Benutzer-Input für die Generierung.

        Returns:
            str: Der Dateipfad der erstellten PDF-Datei.
        """
        content = self.generate_content_with_model("Gemini", user_prompt)
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, content)
        file_path = "erstellte_datei.pdf"
        pdf.output(file_path)
        os.startfile(file_path)  # Datei automatisch öffnen.
        return file_path

    def create_ppt_with_ai(self, user_prompt: str) -> str:
        """
        Erstellt eine PowerPoint-Datei mit von KI generiertem Inhalt.

        Args:
            user_prompt (str): Der Benutzer-Input für die Generierung.

        Returns:
            str: Der Dateipfad der erstellten PowerPoint-Datei.
        """
        content = self.generate_content_with_model("Gemini", user_prompt)
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Erstellte Präsentation"
        subtitle.text = content
        file_path = "erstellte_praesentation.pptx"
        prs.save(file_path)
        os.startfile(file_path)  # Datei automatisch öffnen.
        return file_path

    def create_csv_with_ai(self, user_prompt: str) -> str:
        """
        Erstellt eine CSV-Datei mit von KI generiertem Inhalt.

        Args:
            user_prompt (str): Der Benutzer-Input für die Generierung.

        Returns:
            str: Der Dateipfad der erstellten CSV-Datei.
        """
        content = self.generate_content_with_model("Gemini", user_prompt)
        file_path = "erstellte_datei.csv"
        with open(file_path, mode='w', newline='') as file:
            writer = csv.writer(file)
            rows = content.split("\n")
            for row in rows:
                writer.writerow([col.strip() for col in row.split(",")])
        os.startfile(file_path)  # Datei automatisch öffnen.
        return file_path

# Instanziierung eines FileCreator-Objekts.
file_creator = FileCreator()
