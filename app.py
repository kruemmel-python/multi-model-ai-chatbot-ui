import os
import torch
from openpyxl import Workbook
from docx import Document
from fpdf import FPDF
from pptx import Presentation
import csv
from transformers import AutoModelForSpeechSeq2Seq, AutoProcessor, pipeline
from datasets import load_dataset
import gradio as gr
from PIL import Image
import requests
import json
import base64
from io import BytesIO
from mistralai import Mistral
import google.generativeai as genai
import subprocess
import re
from PyPDF2 import PdfReader
from datetime import datetime
from typing import List, Tuple, Dict, Any, Generator, Optional
from pydub import AudioSegment

# --- Konfiguration ---
device = "cuda:0" if torch.cuda.is_available() else "cpu"
torch_dtype = torch.float16 if torch.cuda.is_available() else torch.float32

model_id = "openai/whisper-large-v3-turbo"

model = AutoModelForSpeechSeq2Seq.from_pretrained(
    model_id, torch_dtype=torch_dtype, low_cpu_mem_usage=True, use_safetensors=True
)
model.to(device)

processor = AutoProcessor.from_pretrained(model_id)

pipe = pipeline(
    "automatic-speech-recognition",
    model=model,
    tokenizer=processor.tokenizer,
    feature_extractor=processor.feature_extractor,
    torch_dtype=torch_dtype,
    device=device,
)

# --- API-Schlüssel und Modelle ---
mistral_api_key = os.environ.get('MISTRAL_API_KEY', 'key')
gemini_api_key = os.environ.get('GEMINI_API_KEY', 'key')

MISTRAL_CHAT_MODEL = "mistral-large-latest"
MISTRAL_IMAGE_MODEL = "pixtral-12b-2409"

MISTRAL_API_URL = "https://api.mistral.ai/v1/chat/completions"

OLLAMA_MODELS = [
    "phi4-model:latest",
    "phi4:latest",
    "wizardlm2:7b-fp16",
    "unzensiert:latest",
    "llama2-uncensored:7b-chat-q8_0",
    "teufel:latest",
    "Odin:latest",
    "luzifer:latest",
    "llama2-uncensored:latest"
]

DEFAULT_OLLAMA_MODEL = "phi4:latest"
STATUS_MESSAGE_GENERATING = "Antwort wird generiert..."
STATUS_MESSAGE_COMPLETE = "Antwort generiert."
STATUS_MESSAGE_ERROR = "Fehler: Die Anfrage konnte nicht verarbeitet werden."

SAVE_DIR = ".gradio"
SAVE_FILE = os.path.join(SAVE_DIR, "save.json")

# --- API Clients Initialisieren ---
mistral_client = Mistral(api_key=mistral_api_key)
genai.configure(api_key=gemini_api_key)

gemini_generation_config = {
    "temperature": 1,
    "top_p": 0.95,
    "top_k": 40,
    "max_output_tokens": 16192,
    "response_mime_type": "text/plain",
}
gemini_model = genai.GenerativeModel(
    model_name="gemini-2.0-flash-exp",
    generation_config=gemini_generation_config,
)

# --- Hilfsfunktionen ---
def encode_image(image: Image.Image) -> str:
    """Kodiert ein Bild in Base64."""
    buffered = BytesIO()
    image.save(buffered, format="JPEG")
    return base64.b64encode(buffered.getvalue()).decode('utf-8')

def format_chat_message(text: str) -> str:
    """Formatiert eine Chatnachricht mit benutzerdefiniertem Stil."""
    return f"<div style='background-color:#333333; padding: 10px; margin-bottom: 5px; border-radius: 5px;'>{text}</div>"

def process_audio(audio_file_path: str) -> str:
    """Verarbeitet eine Audiodatei und extrahiert den Text."""
    try:
        # Audio laden und konvertieren
        audio = AudioSegment.from_file(audio_file_path)
        temp_audio_path = "temp_audio.wav"
        audio.export(temp_audio_path, format="wav")  # Konvertieren in WAV
        # Transkription
        result = pipe(temp_audio_path)
        os.remove(temp_audio_path)  # Temporäre Datei löschen
        return result["text"]
    except Exception as e:
        raise ValueError(f"Fehler bei der Verarbeitung der Audiodatei: {e}")

# --- Funktionen zur Dateierstellung ---
def generate_content_with_model(model_name: str, user_prompt: str) -> str:
    """Generiert den Inhalt basierend auf dem ausgewählten Modell."""
    if model_name == "Mistral":
        response = mistral_client.chat.complete(
            model=MISTRAL_CHAT_MODEL,
            messages=[{"role": "user", "content": user_prompt}]
        )
        return response.choices[0].message.content.strip()
    elif model_name == "Gemini":
        response = gemini_model.generate_content(prompt=user_prompt)
        return response.text.strip()
    elif model_name == "Ollama":
        process = subprocess.Popen(
            ["ollama", "run", DEFAULT_OLLAMA_MODEL],
            stdin=subprocess.PIPE,
            stdout=subprocess.PIPE,
            text=True
        )
        process.stdin.write(user_prompt + "\n")
        process.stdin.close()
        output = process.stdout.read()
        process.stdout.close()
        return clean_output(output).strip()
    else:
        return "Modell nicht verfügbar oder unbekannt."

def create_excel_with_ai(model_name: str, user_prompt: str, sheets: int = 1) -> str:
    """Erstellt eine Excel-Datei mit von KI generiertem Inhalt."""
    content = generate_content_with_model(model_name, user_prompt)
    workbook = Workbook()
    for i in range(sheets):
        sheet = workbook.create_sheet(title=f"Tabelle{i+1}") if i > 0 else workbook.active
        rows = content.split("\n")
        for row_index, row in enumerate(rows, start=1):
            columns = row.split(",")  # Annahme: Kommagetrennte Spalten
            for col_index, value in enumerate(columns, start=1):
                sheet.cell(row=row_index, column=col_index, value=value.strip())
    file_path = "erstellte_tabelle.xlsx"
    workbook.save(file_path)
    os.startfile(file_path)  # Datei automatisch öffnen
    return file_path

def create_word_with_ai(model_name: str, user_prompt: str) -> str:
    """Erstellt eine Word-Datei mit von KI generiertem Inhalt."""
    content = generate_content_with_model(model_name, user_prompt)
    doc = Document()
    doc.add_paragraph(content)
    file_path = "erstelltes_dokument.docx"
    doc.save(file_path)
    os.startfile(file_path)  # Datei automatisch öffnen
    return file_path

def create_pdf_with_ai(model_name: str, user_prompt: str) -> str:
    """Erstellt eine PDF-Datei mit von KI generiertem Inhalt."""
    content = generate_content_with_model(model_name, user_prompt)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, content)
    file_path = "erstellte_datei.pdf"
    pdf.output(file_path)
    os.startfile(file_path)  # Datei automatisch öffnen
    return file_path

def create_ppt_with_ai(model_name: str, user_prompt: str) -> str:
    """Erstellt eine PowerPoint-Datei mit von KI generiertem Inhalt."""
    content = generate_content_with_model(model_name, user_prompt)
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Erstellte Präsentation"
    subtitle.text = content
    file_path = "erstellte_praesentation.pptx"
    prs.save(file_path)
    os.startfile(file_path)  # Datei automatisch öffnen
    return file_path

def create_csv_with_ai(model_name: str, user_prompt: str) -> str:
    """Erstellt eine CSV-Datei mit von KI generiertem Inhalt."""
    content = generate_content_with_model(model_name, user_prompt)
    file_path = "erstellte_datei.csv"
    with open(file_path, mode='w', newline='') as file:
        writer = csv.writer(file)
        rows = content.split("\n")
        for row in rows:
            writer.writerow([col.strip() for col in row.split(",")])
    os.startfile(file_path)  # Datei automatisch öffnen
    return file_path

# --- Mistral Funktionen ---
def chat_with_mistral(user_input: str, chat_history: List[Tuple[str, str]], image: Optional[Image.Image] = None, audio_file: Optional[str] = None) -> Generator[Tuple[List[Tuple[str, str]], str], None, None]:
    if not user_input.strip() and not audio_file:
        yield chat_history, "Bitte geben Sie eine Nachricht ein oder laden Sie eine Audiodatei hoch."
        return

    if audio_file:
        try:
            user_input = process_audio(audio_file)  # Audio in Text umwandeln
        except Exception as e:
            chat_history.append((None, f"Fehler bei der Verarbeitung der Audiodatei: {e}"))
            yield chat_history, ""
            return

    chat_history.append((user_input, None))
    yield chat_history, ""

    messages = [{"role": "user", "content": user_input}]

    if image:
        try:
            image_base64 = encode_image(image)
            messages[0]["content"] = [
                {"type": "text", "text": user_input},
                {"type": "image_url", "image_url": f"data:image/jpeg;base64,{image_base64}"},
            ]
        except Exception as e:
            chat_history.append((None, f"Fehler beim Hochladen des Bildes: {e}"))
            yield chat_history, ""
            return

    try:
        headers = {
            "Content-Type": "application/json",
            "Accept": "application/json",
            "Authorization": f"Bearer {mistral_api_key}",
        }

        payload = {
            "model": MISTRAL_CHAT_MODEL,
            "messages": messages,
            "temperature": 1,
            "top_p": 0.95,
            "max_tokens": 16192,
            "stream": True
        }

        response = requests.post(MISTRAL_API_URL, headers=headers, json=payload, stream=True)
        response.raise_for_status()
        full_response = ""

        for chunk in response.iter_lines():
            if chunk:
                try:
                    # Debug: Logge den ursprünglichen Chunk-Inhalt
                    print(f"Empfangener Chunk: {chunk}")

                    # Ignoriere den speziellen "[DONE]" Chunk
                    if chunk == b"data: [DONE]":
                        print("Stream abgeschlossen (DONE).")
                        break  # Beende die Verarbeitung

                    # Überprüfe, ob der Chunk nicht leer ist
                    if chunk.strip():
                        # Versuche, den Chunk zu dekodieren
                        chunk_data = json.loads(chunk.decode('utf-8').replace('data: ', ''))
                        print(f"Parsed JSON-Daten: {chunk_data}")

                        # Verarbeite den JSON-Inhalt
                        if 'choices' in chunk_data and chunk_data['choices']:
                            delta_content = chunk_data['choices'][0]['delta'].get('content', '')
                            if delta_content:
                                full_response += delta_content
                                formatted_response = format_chat_message(full_response)
                                chat_history[-1] = (user_input, formatted_response)
                                yield chat_history, ""
                except json.JSONDecodeError as e:
                    # Debug: Fehler bei der JSON-Verarbeitung
                    print(f"JSON Decode Fehler: {e} - Ungültiger Chunk: {chunk}")
                    continue
            else:
                # Debug: Leerer Chunk empfangen
                print("Leerer Chunk empfangen.")

        if full_response:
            chat_history[-1] = (user_input, format_chat_message(full_response))
            yield chat_history, ""
        else:
            chat_history.append((None, "Keine Antwort vom Modell erhalten."))
            yield chat_history, ""
    except requests.exceptions.RequestException as e:
        chat_history.append((None, f"Fehler bei der Verarbeitung der Anfrage: {e}"))
        yield chat_history, ""
    except Exception as e:
        chat_history.append((None, f"Unbekannter Fehler: {e}. Bitte versuchen Sie es nochmal."))
        yield chat_history, ""

def analyze_image_mistral(image: Optional[Image.Image], chat_history: List[Tuple[str, str]], user_input: str, prompt: str) -> List[Tuple[str, str]]:
    """Analysiert ein Bild mit Mistral."""
    if image is None:
        chat_history.append((None, "Bitte laden Sie ein Bild hoch."))
        return chat_history

    try:
        image_base64 = encode_image(image)
        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": f"{user_input} {prompt}"},
                    {"type": "image_url", "image_url": f"data:image/jpeg;base64,{image_base64}"},
                ],
            }
        ]

        chat_response = mistral_client.chat.complete(
            model=MISTRAL_IMAGE_MODEL,
            messages=messages
        )

        response_text = chat_response.choices[0].message.content
        chat_history.append((None, format_chat_message(response_text)))
    except Exception as e:
        chat_history.append((None, f"Unbekannter Fehler bei der Bildanalyse: {e}. Bitte versuchen Sie es nochmal."))

    return chat_history

def compare_images_mistral(image1: Optional[Image.Image], image2: Optional[Image.Image], chat_history: List[Tuple[str, str]]) -> List[Tuple[str, str]]:
    """Vergleicht zwei Bilder mit Mistral."""
    if image1 is None or image2 is None:
        chat_history.append((None, "Bitte laden Sie zwei Bilder hoch."))
        return chat_history

    try:
        image1_base64 = encode_image(image1)
        image2_base64 = encode_image(image2)
        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Was sind die Unterschiede zwischen den beiden Bildern? Bitte in Deutsch antworten."},
                    {"type": "image_url", "image_url": f"data:image/jpeg;base64,{image1_base64}"},
                    {"type": "image_url", "image_url": f"data:image/jpeg;base64,{image2_base64}"},
                ],
            }
        ]

        chat_response = mistral_client.chat.complete(
            model=MISTRAL_IMAGE_MODEL,
            messages=messages
        )

        response_text = chat_response.choices[0].message.content
        chat_history.append((None, format_chat_message(response_text)))
    except Exception as e:
        chat_history.append((None, f"Unbekannter Fehler beim Vergleich der Bilder: {e}. Bitte versuchen Sie es nochmal."))

    return chat_history

# --- Gemini Funktionen ---
def upload_to_gemini(image: Image.Image):
    """Lädt ein Bild zur Gemini API hoch."""
    image_path = "temp_image.jpg"
    image.save(image_path)
    sample_file = genai.upload_file(path=image_path, display_name="Hochgeladenes Bild")
    print(f"Hochgeladene Datei '{sample_file.display_name}' as: {sample_file.uri}")
    os.remove(image_path)
    return sample_file

def chat_with_gemini(
    user_input: str,
    chat_history: List[Tuple[str, str]],
    image: Optional[Image.Image] = None,
    audio_file: Optional[str] = None
) -> Generator[Tuple[List[Tuple[str, str]], str], None, None]:
    if not user_input.strip() and not audio_file:
        yield chat_history, "Bitte geben Sie eine Nachricht ein oder laden Sie eine Audiodatei hoch."
        return

    if audio_file:
        try:
            user_input = process_audio(audio_file)  # Audio in Text umwandeln
        except Exception as e:
            chat_history.append((None, f"Fehler bei der Verarbeitung der Audiodatei: {e}"))
            yield chat_history, ""
            return

    chat_history.append((user_input, None))
    yield chat_history, ""

    # Chatverlauf initialisieren
    history = [{"role": "user", "parts": [user_input]}]

    if image:
        try:
            sample_file = upload_to_gemini(image)
            history[0]["parts"].append(sample_file)
        except Exception as e:
            chat_history.append((None, f"Fehler beim Hochladen des Bildes: {e}"))
            yield chat_history, ""
            return

    try:
        # Chat mit Gemini starten
        chat_session = gemini_model.start_chat(history=history)
        response_text = chat_session.send_message(user_input).text
        chat_history.append((None, format_chat_message(response_text)))
    except Exception as e:
        chat_history.append((None, f"Fehler bei der Verarbeitung der Anfrage: {e}"))

    yield chat_history, ""

def analyze_image_gemini(image: Optional[Image.Image], chat_history: List[Tuple[str, str]], user_input: str) -> List[Tuple[str, str]]:
    """Analysiert ein Bild mit Gemini."""
    if image is None:
        chat_history.append((None, "Bitte laden Sie ein Bild hoch."))
        return chat_history
    try:
        sample_file = upload_to_gemini(image)
        response = gemini_model.generate_content([f"{user_input} Beschreiben Sie das Bild mit einer kreativen Beschreibung. Bitte in German antworten.", sample_file])
        response_text = response.text
        chat_history.append((None, format_chat_message(response_text)))
    except Exception as e:
        chat_history.append((None, f"Fehler bei der Bildanalyse: {e}"))

    return chat_history

# --- Chat-Verwaltungsfunktionen ---
def clear_chat(chat_history: List[Tuple[str, str]]) -> List[Tuple[str, str]]:
    """Leert den Chatverlauf."""
    chat_history.clear()
    return chat_history

def generate_chat_title(chat_history: List[Tuple[str, str]]) -> str:
    """Generiert einen Titel basierend auf den ersten Chatnachrichten."""
    if not chat_history:
        return "Neuer Chat"
    first_messages = [msg[0] for msg in chat_history if msg[0] is not None][:2]
    if not first_messages:
        return "Neuer Chat"
    title = " ".join(first_messages)
    title = title[:50] + "..." if len(title) > 50 else title
    return title

def save_chat(chat_history: List[Tuple[str, str]], saved_chats: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Speichert den aktuellen Chatverlauf."""
    if not chat_history:
        return saved_chats

    title = generate_chat_title(chat_history)
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    saved_chat = {
        "title": title,
        "date": timestamp,
        "chat": chat_history,
    }
    saved_chats.append(saved_chat)
    _save_chats_to_file(saved_chats)
    return saved_chats

def _save_chats_to_file(chats: List[Dict[str, Any]]):
    """Speichert Chatverläufe in einer JSON-Datei."""
    os.makedirs(SAVE_DIR, exist_ok=True)
    with open(SAVE_FILE, 'w', encoding='utf-8') as f:
        json.dump(chats, f, ensure_ascii=False, indent=4)

def _load_chats_from_file() -> List[Dict[str, Any]]:
    """Lädt Chatverläufe aus einer JSON-Datei."""
    if os.path.exists(SAVE_FILE):
        try:
            with open(SAVE_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except json.JSONDecodeError:
            return []
    return []

def format_saved_chat(saved_chat: Dict[str, Any]) -> str:
    """Formatiert einen gespeicherten Chat für die Anzeige."""
    return f"**{saved_chat['title']}** ({saved_chat['date']})"

def load_chat(chat_title: str, saved_chats: List[Dict[str, Any]], chat_history: List[Tuple[str, str]]) -> Tuple[List[Tuple[str, str]], List[Tuple[str, str]], str]:
    """Lädt einen Chat in die Chat-Ausgabe."""
    selected_chat = next((chat for chat in saved_chats if format_saved_chat(chat) == chat_title), None)
    if selected_chat:
        return selected_chat['chat'], selected_chat['chat'], chat_title
    return chat_history, chat_history, None  # Fallback if title not found

def new_chat(saved_chats: List[Dict[str, Any]]) -> Tuple[List[Tuple[str, str]], List[Dict[str, Any]]]:
    """Startet einen neuen Chat."""
    return [], saved_chats

def delete_chat(chat_title: str, saved_chats: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Löscht einen einzelnen Chat."""
    updated_chats = [chat for chat in saved_chats if format_saved_chat(chat) != chat_title]
    _save_chats_to_file(updated_chats)
    return updated_chats

def delete_all_chats() -> List[Dict[str, Any]]:
    """Löscht alle gespeicherten Chats."""
    _save_chats_to_file([])
    return []

# --- Ollama Funktionen ---
def clean_output(output: str) -> str:
    """Entfernt Steuerzeichen aus der Ollama-Ausgabe."""
    cleaned_output = re.sub(r'(?:\x1B[@-_]|[\x1B\x9B][0-?]*[ -/]*[@-~])', '', output)
    cleaned_output = re.sub(r'\?\d+[lh]', '', cleaned_output)
    cleaned_output = re.sub(r'[\u2800-\u28FF]', '', cleaned_output)
    cleaned_output = re.sub(r'\r', '', cleaned_output)
    cleaned_output = re.sub(r'2K1G ?(?:2K1G)*!?', '', cleaned_output)
    return cleaned_output

def format_as_codeblock(output: str) -> str:
    """Formatiert die Ausgabe als Markdown-Codeblock."""
    return f"```\n{output}\n```"

def run_ollama_live(prompt: str, model: str) -> Generator[str, None, None]:
    """Führt Ollama aus und gibt die Ausgabe live zurück."""
    try:
        process = subprocess.Popen(
            ["ollama", "run", model],
            stdin=subprocess.PIPE,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            bufsize=1,
            encoding='utf-8',
        )
        process.stdin.write(prompt + "\n")
        process.stdin.close()

        buffer = ""
        for line in iter(process.stdout.readline, ''):
            clean_line = clean_output(line)
            if clean_line:
                buffer += clean_line + "\n"
                yield format_as_codeblock(buffer)

        process.stdout.close()
        process.wait()

    except Exception as e:
        yield f"**Fehler:** {str(e)}"

def process_uploaded_file(file: gr.File) -> str:
    """Verarbeitet hochgeladene TXT- und PDF-Dateien."""
    if file.name.endswith(".txt"):
        with open(file.name, 'r', encoding='utf-8') as f:
            return f.read()
    elif file.name.endswith(".pdf"):
        try:
            reader = PdfReader(file.name)
            content = ""
            for page in reader.pages:
                content += page.extract_text()
            return content
        except Exception as e:
            raise ValueError(f"Fehler beim Lesen der PDF-Datei: {e}")
    else:
        raise ValueError("Nur TXT- und PDF-Dateien werden unterstützt.")

def chatbot_interface(input_text: str, model: str, file: Optional[gr.File], audio_file: Optional[gr.File] = None) -> Generator[Tuple[str, str], None, None]:
    """Schnittstelle für die Ollama-Chatbot-Funktion."""
    yield "", STATUS_MESSAGE_GENERATING

    if file:
        try:
            input_text = process_uploaded_file(file)
        except ValueError as e:
            yield f"**Fehler:** {str(e)}", STATUS_MESSAGE_ERROR
            return
        except Exception as e:
            yield f"**Unerwarteter Fehler:** {str(e)}", STATUS_MESSAGE_ERROR
            return

    if audio_file:
        try:
            input_text = process_audio(audio_file)  # Übergibt den Dateipfad direkt
        except Exception as e:
            yield f"**Fehler bei der Verarbeitung der Audiodatei:** {str(e)}", STATUS_MESSAGE_ERROR
            return

    try:
        for chunk in run_ollama_live(input_text, model):
            yield chunk, STATUS_MESSAGE_GENERATING
        yield chunk, STATUS_MESSAGE_COMPLETE
    except Exception as e:
        yield f"**Fehler bei der Kommunikation mit Ollama:** {str(e)}", STATUS_MESSAGE_ERROR

# --- Gradio Benutzeroberfläche ---
with gr.Blocks() as demo:
    gr.Markdown("## Chatbots mit Bild- und Audioanalyse")

    with gr.Tabs():
        # --- Mistral Chatbot ---
        with gr.TabItem("Mistral Chatbot"):
            mistral_chatbot = gr.Chatbot(label="Chat-Verlauf", height=400)
            mistral_state = gr.State([])
            mistral_saved_chats = gr.State(_load_chats_from_file())

            with gr.Row():
                with gr.Column(scale=5):
                    mistral_user_input = gr.Textbox(label="Nachricht", placeholder="Geben Sie hier Ihre Nachricht ein...")
                with gr.Column(scale=1, min_width=100):
                    mistral_submit_btn = gr.Button("Senden")
            with gr.Row():
                with gr.Column(scale=1):
                    mistral_image_upload = gr.Image(type="pil", label="Bild hochladen", height=200)
                with gr.Column(scale=1):
                    mistral_audio_upload = gr.Audio(type="filepath", label="Audio hochladen")  # NEU
                with gr.Column(scale=1):
                    mistral_analyze_btn = gr.Button("Bild Nachricht senden")

            with gr.Row():
                mistral_clear_chat_button = gr.Button("Chat leeren")
                mistral_save_chat_button = gr.Button("Chat speichern")
                mistral_new_chat_button = gr.Button("Neuen Chat starten")

            with gr.Accordion("Gespeicherte Chats", open=False):
                mistral_saved_chat_display = gr.Radio(label="Gespeicherte Chats", interactive=True)
                with gr.Row():
                    mistral_delete_chat_button = gr.Button("Ausgewählten Chat löschen")
                    mistral_delete_all_chats_button = gr.Button("Alle Chats löschen")

            with gr.Tabs():
                with gr.TabItem("Bildanalyse"):
                    with gr.Row():
                        with gr.Column(scale=1):
                            mistral_image_upload_analysis = gr.Image(type="pil", label="Bild hochladen", height=200)
                        with gr.Column(scale=1):
                            mistral_analyze_btn_analysis = gr.Button("Nur Bild analysieren")

                    with gr.Row():
                        with gr.Column(scale=1):
                            mistral_analyze_btn_charts = gr.Button("Diagramme verstehen")
                        with gr.Column(scale=1):
                            mistral_analyze_btn_compare = gr.Button("Bilder vergleichen")
                        with gr.Column(scale=1):
                            mistral_analyze_btn_receipts = gr.Button("Belege umschreiben")
                        with gr.Column(scale=1):
                            mistral_analyze_btn_documents = gr.Button("Alte Dokumente umschreiben")
                        with gr.Column(scale=1):
                            mistral_analyze_btn_ocr = gr.Button("OCR mit strukturiertem Output")

                with gr.TabItem("Bildvergleich"):
                    with gr.Row():
                        with gr.Column(scale=1):
                            mistral_image_upload_compare1 = gr.Image(type="pil", label="Bild 1 hochladen", height=200)
                        with gr.Column(scale=1):
                            mistral_image_upload_compare2 = gr.Image(type="pil", label="Bild 2 hochladen", height=200)
                    with gr.Row():
                        mistral_compare_btn = gr.Button("Bilder vergleichen")

            mistral_submit_btn.click(
                chat_with_mistral,
                inputs=[mistral_user_input, mistral_state, mistral_image_upload, mistral_audio_upload],  # Audio hinzugefügt
                outputs=[mistral_chatbot, mistral_user_input]
            )

            mistral_analyze_btn.click(
                lambda image, chat_history, user_input: analyze_image_mistral(image, chat_history, user_input, "Beschreiben Sie das Bild mit einer kreativen Beschreibung. Bitte in Deutsch antworten."),
                inputs=[mistral_image_upload, mistral_state, mistral_user_input],
                outputs=[mistral_chatbot]
            )
            mistral_clear_chat_button.click(clear_chat, inputs=[mistral_state], outputs=[mistral_chatbot])
            mistral_save_chat_button.click(save_chat, inputs=[mistral_state, mistral_saved_chats], outputs=[mistral_saved_chats])
            mistral_new_chat_button.click(new_chat, inputs=[mistral_saved_chats], outputs=[mistral_chatbot, mistral_saved_chats])

            def update_mistral_radio(chats):
                formatted_chats = [format_saved_chat(chat) for chat in chats]
                return gr.update(choices=formatted_chats)

            mistral_saved_chats.change(update_mistral_radio, inputs=[mistral_saved_chats], outputs=[mistral_saved_chat_display])
            mistral_saved_chat_display.change(load_chat, inputs=[mistral_saved_chat_display, mistral_saved_chats, mistral_state], outputs=[mistral_chatbot, mistral_state, mistral_saved_chat_display])

            mistral_delete_chat_button.click(delete_chat, inputs=[mistral_saved_chat_display, mistral_saved_chats], outputs=[mistral_saved_chats])
            mistral_delete_all_chats_button.click(delete_all_chats, outputs=[mistral_saved_chats])
            demo.load(update_mistral_radio, inputs=[mistral_saved_chats], outputs=[mistral_saved_chat_display])

            mistral_analyze_btn_analysis.click(
                lambda image, chat_history, user_input: analyze_image_mistral(image, chat_history, user_input, "Was ist auf diesem Bild? Bitte in Deutsch antworten."),
                inputs=[mistral_image_upload_analysis, mistral_state, mistral_user_input],
                outputs=[mistral_chatbot]
            )
            mistral_analyze_btn_charts.click(
                lambda image, chat_history, user_input: analyze_image_mistral(image, chat_history, user_input, "Was ist auf diesem Bild? Bitte in Deutsch antworten."),
                inputs=[mistral_image_upload_analysis, mistral_state, mistral_user_input],
                outputs=[mistral_chatbot]
            )
            mistral_analyze_btn_compare.click(
                lambda image, chat_history, user_input: analyze_image_mistral(image, chat_history, user_input, "Was sind die Unterschiede zwischen den beiden Bildern? Bitte in Deutsch antworten."),
                inputs=[mistral_image_upload_analysis, mistral_state, mistral_user_input],
                outputs=[mistral_chatbot]
            )
            mistral_analyze_btn_receipts.click(
                lambda image, chat_history, user_input: analyze_image_mistral(image, chat_history, user_input, "Transkribieren Sie diesen Beleg. Bitte in Deutsch antworten."),
                inputs=[mistral_image_upload_analysis, mistral_state, mistral_user_input],
                outputs=[mistral_chatbot]
            )
            mistral_analyze_btn_documents.click(
                lambda image, chat_history, user_input: analyze_image_mistral(image, chat_history, user_input, "Transkribieren Sie dieses Dokument. Bitte in Deutsch antworten."),
                inputs=[mistral_image_upload_analysis, mistral_state, mistral_user_input],
                outputs=[mistral_chatbot]
            )
            mistral_analyze_btn_ocr.click(
                lambda image, chat_history, user_input: analyze_image_mistral(image, chat_history, user_input, "Extrahieren Sie aus dieser Rechnung die Rechnungsnummer, die Artikelnamen und zugehörigen Preise sowie den Gesamtpreis und geben Sie sie als Zeichenfolge in einem JSON-Objekt zurück. Bitte in Deutsch antworten."),
                inputs=[mistral_image_upload_analysis, mistral_state, mistral_user_input],
                outputs=[mistral_chatbot]
            )

            mistral_compare_btn.click(compare_images_mistral, inputs=[mistral_image_upload_compare1, mistral_image_upload_compare2, mistral_state], outputs=[mistral_chatbot])

        # --- Gemini Chatbot ---
        with gr.TabItem("Gemini Chatbot"):
            gemini_chatbot = gr.Chatbot(label="Chat-Verlauf", height=400)
            gemini_state = gr.State([])
            gemini_saved_chats = gr.State(_load_chats_from_file())

            with gr.Row():
                with gr.Column(scale=5):
                    gemini_user_input = gr.Textbox(label="Nachricht", placeholder="Geben Sie hier Ihre Nachricht ein...")
                with gr.Column(scale=1, min_width=100):
                    gemini_submit_btn = gr.Button("Senden")
            with gr.Row():
                with gr.Column(scale=1):
                    gemini_image_upload = gr.Image(type="pil", label="Bild hochladen", height=200)
                with gr.Column(scale=1):
                    gemini_audio_upload = gr.Audio(type="filepath", label="Audio hochladen")  # NEU
                with gr.Column(scale=1):
                    gemini_analyze_btn = gr.Button("Bild mit Nachricht senden")

            with gr.Row():
                gemini_clear_chat_button = gr.Button("Chat leeren")
                gemini_save_chat_button = gr.Button("Chat speichern")
                gemini_new_chat_button = gr.Button("Neuen Chat starten")

            with gr.Accordion("Gespeicherte Chats", open=False):
                gemini_saved_chat_display = gr.Radio(label="Gespeicherte Chats", interactive=True)
                with gr.Row():
                    gemini_delete_chat_button = gr.Button("Ausgewählten Chat löschen")
                    gemini_delete_all_chats_button = gr.Button("Alle Chats löschen")

            gemini_submit_btn.click(
                chat_with_gemini,
                inputs=[gemini_user_input, gemini_state, gemini_image_upload, gemini_audio_upload],  # Audio hinzugefügt
                outputs=[gemini_chatbot, gemini_user_input]
            )

            gemini_analyze_btn.click(
                lambda image, chat_history, user_input: analyze_image_gemini(image, chat_history, user_input),
                inputs=[gemini_image_upload, gemini_state, gemini_user_input],
                outputs=[gemini_chatbot]
            )
            gemini_clear_chat_button.click(clear_chat, inputs=[gemini_state], outputs=[gemini_chatbot])
            gemini_save_chat_button.click(save_chat, inputs=[gemini_state, gemini_saved_chats], outputs=[gemini_saved_chats])
            gemini_new_chat_button.click(new_chat, inputs=[gemini_saved_chats], outputs=[gemini_chatbot, gemini_saved_chats])

            def update_gemini_radio(chats):
                formatted_chats = [format_saved_chat(chat) for chat in chats]
                return gr.update(choices=formatted_chats)

            gemini_saved_chats.change(update_gemini_radio, inputs=[gemini_saved_chats], outputs=[gemini_saved_chat_display])
            gemini_saved_chat_display.change(load_chat, inputs=[gemini_saved_chat_display, gemini_saved_chats, gemini_state], outputs=[gemini_chatbot, gemini_state, gemini_saved_chat_display])

            gemini_delete_chat_button.click(delete_chat, inputs=[gemini_saved_chat_display, gemini_saved_chats], outputs=[gemini_saved_chats])
            gemini_delete_all_chats_button.click(delete_all_chats, outputs=[gemini_saved_chats])
            demo.load(update_gemini_radio, inputs=[gemini_saved_chats], outputs=[gemini_saved_chat_display])

        # --- Ollama Chatbot ---
        with gr.TabItem("Ollama Chatbot"):
            ollama_input_text = gr.Textbox(lines=2, placeholder="Geben Sie Ihre Frage an Ollama ein", label="Eingabe (oder Datei hochladen)")
            ollama_model_selector = gr.Dropdown(choices=OLLAMA_MODELS, value=DEFAULT_OLLAMA_MODEL, label="Modell auswählen")
            ollama_file_upload = gr.File(label="Datei hochladen (PDF oder TXT)", file_types=[".txt", ".pdf"])
            ollama_audio_upload = gr.Audio(type="filepath", label="Audio hochladen")
            ollama_output = gr.Markdown(label="Antwort")
            ollama_status = gr.Label(label="Status")

            ollama_submit_btn = gr.Button("Senden")
            ollama_submit_btn.click(chatbot_interface, inputs=[ollama_input_text, ollama_model_selector, ollama_file_upload, ollama_audio_upload], outputs=[ollama_output, ollama_status])

        # --- Dateierstellung ---
        with gr.TabItem("Dateierstellung"):
            gr.Markdown("## Dateierstellung")
            with gr.Row():
                with gr.Column(scale=1):
                    file_type = gr.Radio(choices=["Excel", "Word", "PDF", "PowerPoint", "CSV"], label="Dateiformat", value="Excel")
                with gr.Column(scale=1):
                    sheets = gr.Slider(minimum=1, maximum=5, step=1, label="Anzahl der Tabellenblätter (nur Excel)", value=1)
                with gr.Column(scale=2):
                    file_content = gr.Textbox(label="Inhalt der Datei", placeholder="Geben Sie den Inhalt der Datei ein...")
                with gr.Column(scale=1):
                    create_file_btn = gr.Button("Datei erstellen")
                    download_file = gr.File(label="Herunterladen")

            create_file_btn.click(
                lambda content, file_format, sheets: (
                    create_excel_with_ai("Mistral", content, sheets) if file_format == "Excel" else
                    create_word_with_ai("Mistral", content) if file_format == "Word" else
                    create_pdf_with_ai("Mistral", content) if file_format == "PDF" else
                    create_ppt_with_ai("Mistral", content) if file_format == "PowerPoint" else
                    create_csv_with_ai("Mistral", content)
                ),
                inputs=[file_content, file_type, sheets],
                outputs=[download_file]
            )

# Starten Sie die Gradio-Benutzeroberfläche
if __name__ == '__main__':
    demo.launch(share=True, server_name="localhost", server_port=8779)
